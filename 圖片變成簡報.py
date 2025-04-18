# 將圖片轉為PPT的 Streamlit 版本（介面與功能與原始版本一致）

import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import os
import io
import zipfile

def convert_images_to_ppt(uploaded_files, resize_option, width, height):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    for uploaded_file in uploaded_files:
        image = Image.open(uploaded_file)
        filename = os.path.splitext(uploaded_file.name)[0]

        slide = prs.slides.add_slide(blank_slide_layout)

        if resize_option == "resize":
            image = image.resize((width, height))
            temp_bytes = io.BytesIO()
            image.save(temp_bytes, format="PNG")
            temp_bytes.seek(0)
            slide.shapes.add_picture(temp_bytes, Inches(1), Inches(1), width=Inches(8))
        else:
            temp_bytes = io.BytesIO()
            image.save(temp_bytes, format="PNG")
            temp_bytes.seek(0)
            slide.shapes.add_picture(temp_bytes, Inches(1), Inches(1), width=Inches(8))

        # 標題
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.2), Inches(8), Inches(1))
        tf = title_box.text_frame
        tf.text = filename
        tf.paragraphs[0].font.size = Pt(24)

    ppt_bytes = io.BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    return ppt_bytes

# 建立 Streamlit 網頁
st.title("🖼️ 圖片轉 PPT 工具")
st.markdown("將多張圖片一鍵轉成 PowerPoint，每張圖片一頁，並用檔名作為標題。")

uploaded_files = st.file_uploader("請選擇圖片（可多選）", type=["png", "jpg", "jpeg", "bmp", "gif"], accept_multiple_files=True)

resize_option = st.radio("圖片處理方式：", ["keep", "resize"], format_func=lambda x: "維持原尺寸" if x=="keep" else "縮放成指定尺寸")

col1, col2 = st.columns(2)
with col1:
    width = st.number_input("寬度 (px)", value=800, step=50)
with col2:
    height = st.number_input("高度 (px)", value=600, step=50)

if st.button("🚀 產生 PPT"):
    if uploaded_files:
        ppt_data = convert_images_to_ppt(uploaded_files, resize_option, int(width), int(height))
        st.success("✅ 轉換完成！")
        st.download_button("⬇️ 下載 PPT", ppt_data, file_name="converted.pptx")
    else:
        st.warning("請先上傳圖片喔！")
