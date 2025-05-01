import streamlit as st
from pptx import Presentation
from pptx.util import Cm, Pt
from PIL import Image
import io
import datetime

st.set_page_config(page_title="圖轉PPT", page_icon="🐰", layout="centered")
st.title("🐰 圖轉PPT小工具")

# Sidebar 設定
with st.sidebar:
    layout = st.selectbox("選擇排版樣式", ["1圖一頁", "5圖一頁", "8圖一頁"], index=2)
    margin = st.number_input("簡報邊界（cm）", min_value=0.5, value=1.0, step=0.5)
    padding = st.number_input("圖片間距（px）", min_value=0, value=2, step=1)
    quality_option = st.selectbox("圖片壓縮品質", ["原圖（無壓縮）", "建議（輕壓縮）", "小檔（高壓縮）"], index=1)
    quality_dict = {"原圖（無壓縮）": 100, "建議（輕壓縮）": 85, "小檔（高壓縮）": 65}

# 初始化 session_state
if "images" not in st.session_state:
    st.session_state.images = []

# 上傳圖片區塊
uploaded = st.file_uploader(
    "📤 拖曳圖片到這裡（可多次補充上傳，最多 200MB）",
    type=["png", "jpg", "jpeg", "bmp", "gif"],
    accept_multiple_files=True,
    key="uploader"
)

# 插圖共用函式
def insert_image(slide, image_file, center_pos, max_width, max_height, quality):
    img = Image.open(image_file)
    img_ratio = img.width / img.height
    box_ratio = max_width / max_height
    if img_ratio > box_ratio:
        final_width = max_width
        final_height = max_width / img_ratio
    else:
        final_height = max_height
        final_width = max_height * img_ratio
    buffer = io.BytesIO()
    img.convert("RGB").save(buffer, format='JPEG', quality=quality)
    buffer.seek(0)
    pos_x, pos_y = center_pos
    slide.shapes.add_picture(buffer, pos_x - final_width / 2, pos_y - final_height / 2, width=final_width, height=final_height)


# 加入新上傳圖片（避免重複）
if uploaded:
    for file in uploaded:
        if file not in st.session_state.images:
            st.session_state.images.append(file)

# 顯示圖片數與操作按鈕
if st.session_state.images:
    st.success(f"✅ 已上傳 {len(st.session_state.images)} 張圖片")

    buffer = None
#if st.session_state.images:
    now = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"圖片簡報_{now}.pptx"

    prs = Presentation()
    if layout == "1圖一頁":
        for image in st.session_state.images:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            w, h = prs.slide_width, prs.slide_height
            margin_cm = Cm(margin)
            usable_w = w - 2 * margin_cm
            usable_h = h - 2 * margin_cm
            insert_image(slide, image, (w/2, h/2), usable_w, usable_h, quality_dict[quality_option])

    elif layout == "5圖一頁":
        layout_order = ["center", "top", "bottom", "left", "right"]
        for i in range(0, len(st.session_state.images), 5):
            subset = st.session_state.images[i:i+5]
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            w, h = prs.slide_width, prs.slide_height
            m = Cm(margin)
            px = Pt(padding)
            usable_w = w - 2 * m
            usable_h = h - 2 * m
            cx, cy = w / 2, h / 2
            positions = {
                "center": (cx, cy),
                "top": (cx, px + cy - usable_h / 3),
                "bottom": (cx, cy - px + usable_h / 3),
                "left": (cx - px - usable_w / 3, cy),
                "right": (px + cx + usable_w / 3, cy)
            }
            max_w = usable_w / 3
            max_h = usable_h / 3
            for j, image in enumerate(subset):
                if j >= len(layout_order):
                    break
                insert_image(slide, image, positions[layout_order[j]], max_w, max_h, quality_dict[quality_option])

    elif layout == "8圖一頁":
        nine_grid_order = [4, 6, 0, 2, 3, 5, 7, 1]
        for i in range(0, len(st.session_state.images), 8):
            subset = st.session_state.images[i:i+8]
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            w, h = prs.slide_width, prs.slide_height
            m = Cm(margin)
            usable_w = w - 2 * m
            usable_h = h - 2 * m
            cols, rows = 3, 3
            cell_w = usable_w / cols
            cell_h = usable_h / rows
            px = Pt(padding)
            max_w = cell_w - px
            max_h = cell_h - px
            for draw_idx in range(len(subset)):
                pos = nine_grid_order[draw_idx]
                row = pos // cols
                col = pos % cols
                pos_x = m + col * cell_w + cell_w / 2
                pos_y = m + row * cell_h + cell_h / 2
                insert_image(slide, subset[draw_idx], (pos_x, pos_y), max_w, max_h, quality_dict[quality_option])

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    st.download_button(
        "🚀 產生並下載簡報",
        data=buffer,
        file_name=filename,
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

