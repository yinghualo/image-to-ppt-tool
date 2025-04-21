from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
from PIL import Image

# 圖片資料夾路徑（可以改成自己的）
img_folder = Path("C:/Users/chihane/Desktop/拍賣")  # 例如 Path("C:/Users/你/Downloads/images")

# 建立簡報
prs = Presentation()

# 頁面尺寸（這是預設尺寸，可依需求改）
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# 支援圖片副檔名
img_exts = [".jpg", ".jpeg", ".png"]

# 一張圖一頁
for img_path in sorted(img_folder.glob("*")):
    if img_path.suffix.lower() in img_exts:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # 空白頁面

        # 標題（圖檔名稱不含副檔名）
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = img_path.stem
        title_frame.paragraphs[0].font.size = Pt(24)

        # 插入圖片，讓圖片盡量填滿頁面（不超過邊界）
        img = Image.open(img_path)
        img_width, img_height = img.size
        aspect_ratio = img_width / img_height

        max_width = prs.slide_width - Inches(1)
        max_height = prs.slide_height - Inches(2)

        if aspect_ratio > max_width / max_height:
            width = max_width
            height = width / aspect_ratio
        else:
            height = max_height
            width = height * aspect_ratio

        left = (prs.slide_width - width) / 2
        top = Inches(1.2)

        slide.shapes.add_picture(str(img_path), left, top, width=width, height=height)

# 存檔
prs.save("output.pptx")
print("✅ 簡報已產生！")