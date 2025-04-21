import os
from tkinter import Tk, filedialog, Label, Button, Radiobutton, IntVar, Entry
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

def convert_images_to_ppt(folder_path, save_path, resize_option, width=800, height=600):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    for filename in os.listdir(folder_path):
        if filename.lower().endswith((".png", ".jpg", ".jpeg", ".bmp", ".gif")):
            img_path = os.path.join(folder_path, filename)
            slide = prs.slides.add_slide(blank_slide_layout)

            image = Image.open(img_path)
            if resize_option == 2:
                image = image.resize((int(width), int(height)))
                temp_path = os.path.join(folder_path, "resized_temp_image.png")
                image.save(temp_path)
                img_to_use = temp_path
            else:
                img_to_use = img_path

            slide.shapes.add_picture(img_to_use, Inches(1), Inches(1), width=Inches(8))

            title_box = slide.shapes.add_textbox(Inches(1), Inches(0.2), Inches(8), Inches(1))
            tf = title_box.text_frame
            tf.text = os.path.splitext(filename)[0]
            tf.paragraphs[0].font.size = Pt(24)

    prs.save(save_path)
    print(f"✅ 已儲存至：{save_path}")

# ====== UI ======
def browse_folder():
    folder = filedialog.askdirectory()
    folder_entry.delete(0, 'end')
    folder_entry.insert(0, folder)

def browse_save_path():
    path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint", "*.pptx")])
    save_entry.delete(0, 'end')
    save_entry.insert(0, path)

def start_conversion():
    folder = folder_entry.get()
    save_path = save_entry.get()
    option = resize_var.get()
    w = width_entry.get()
    h = height_entry.get()

    if option == 2 and w and h:
        convert_images_to_ppt(folder, save_path, option, int(w), int(h))
    else:
        convert_images_to_ppt(folder, save_path, option)
    status_label.config(text="✅ 轉換完成！")

root = Tk()
root.title("圖片轉PPT工具")

Label(root, text="📁 圖片資料夾：").grid(row=0, column=0, sticky='w')
folder_entry = Entry(root, width=40)
folder_entry.grid(row=0, column=1)
Button(root, text="選擇", command=browse_folder).grid(row=0, column=2)

Label(root, text="💾 儲存位置：").grid(row=1, column=0, sticky='w')
save_entry = Entry(root, width=40)
save_entry.grid(row=1, column=1)
Button(root, text="選擇", command=browse_save_path).grid(row=1, column=2)

resize_var = IntVar(value=1)
Label(root, text="📏 圖片處理方式：").grid(row=2, column=0, sticky='w')
Radiobutton(root, text="維持原比例", variable=resize_var, value=1).grid(row=2, column=1, sticky='w')
Radiobutton(root, text="縮放為：", variable=resize_var, value=2).grid(row=3, column=1, sticky='w')

width_entry = Entry(root, width=5)
width_entry.insert(0, "800")
width_entry.grid(row=3, column=1, padx=(80, 0), sticky='w')

Label(root, text="×").grid(row=3, column=1, padx=(130, 0), sticky='w')

height_entry = Entry(root, width=5)
height_entry.insert(0, "600")
height_entry.grid(row=3, column=1, padx=(150, 0), sticky='w')

Button(root, text="🚀 開始轉換", command=start_conversion).grid(row=4, column=1, pady=10)
status_label = Label(root, text="")
status_label.grid(row=5, column=1)

root.mainloop()
